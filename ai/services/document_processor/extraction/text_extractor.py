import json
import logging
import re as _re
from pathlib import Path
from typing import FrozenSet, Optional
import asyncio

from .file_validator import BASE_SUPPORTED_EXTENSIONS, IMAGE_EXTENSIONS_OCR

logger = logging.getLogger(__name__)

# No built-in document types — the application is domain-agnostic by default.
# Populate ``ai/resources/document_category_taxonomy.json`` with your domain's
# document type names (e.g. "INVOICE", "CONTRACT", "REPORT") so that
# extract_doc_title() and the listing classifier can recognise them.
_DEFAULT_STANDALONE_DOC_TYPES: FrozenSet[str] = frozenset()

# ``ai/services/document_processor/extraction`` → ``ai/resources``
_TAXONOMY_JSON = Path(__file__).resolve().parents[3] / "resources" / "document_category_taxonomy.json"


def _load_standalone_doc_types() -> FrozenSet[str]:
    """Load standalone-line labels from taxonomy JSON. Returns empty set if not configured."""
    if not _TAXONOMY_JSON.is_file():
        return _DEFAULT_STANDALONE_DOC_TYPES
    try:
        with open(_TAXONOMY_JSON, encoding="utf-8") as f:
            data = json.load(f)
        raw = data.get("standalone_types")
        if not isinstance(raw, list):
            return _DEFAULT_STANDALONE_DOC_TYPES
        return frozenset(str(t).strip().upper() for t in raw if str(t).strip())
    except OSError as e:
        logger.warning("Could not read document taxonomy at %s: %s", _TAXONOMY_JSON, e)
        return _DEFAULT_STANDALONE_DOC_TYPES
    except json.JSONDecodeError as e:
        logger.warning("Invalid JSON in document taxonomy %s: %s", _TAXONOMY_JSON, e)
        return _DEFAULT_STANDALONE_DOC_TYPES


# Known document-type labels (standalone line). Prefer deployment JSON when present.
STANDALONE_DOC_TYPES: FrozenSet[str] = _load_standalone_doc_types()

# How far into the document to scan for a type label.
_PRIMARY_TYPE_SCAN_CHARS = 16000


def _normalize_title_line(line: str) -> str:
    """Collapse internal whitespace for robust match against OCR / large headings."""
    return _re.sub(r"\s+", " ", line.strip()).upper()


def extract_doc_title(text: str) -> Optional[str]:
    """
    Return the first line of *text* that exactly matches a configured STANDALONE_DOC_TYPES
    label (after whitespace normalization). Returns None when no match is found or when the
    taxonomy is empty (the default — see ai/resources/document_category_taxonomy.json).
    """
    if not text or not STANDALONE_DOC_TYPES:
        return None
    for line in text[:_PRIMARY_TYPE_SCAN_CHARS].splitlines():
        u = _normalize_title_line(line)
        if u in STANDALONE_DOC_TYPES:
            return u
    return None


def build_layout_aware_preview(text: str, max_chars: int = 500) -> str:
    """
    Build a stable preview from layout-marked extracted text.

    When PDF extraction includes [Region: ...] markers, taking the first N chars
    can hide the right-side panel (e.g. REFERENCES) behind long left/body text.
    This composes a preview by taking small slices from each region in a fixed
    order for Page 1: full → left → right.
    """
    if not text:
        return ""
    if "[Region:" not in text or "[Page 1]" not in text:
        return (text[:max_chars] if len(text) > max_chars else text).strip()

    # Extract the first page only for preview composition.
    page1 = text.split("[Page 1]", 1)[1]
    if "\n\n[Page " in page1:
        page1 = page1.split("\n\n[Page ", 1)[0]

    def _grab(region: str) -> str:
        token = f"[Region: {region}]"
        if token not in page1:
            return ""
        seg = page1.split(token, 1)[1]
        seg = seg.split("[Region:", 1)[0]
        return seg.strip()

    full = _grab("full")
    left = _grab("left")
    right = _grab("right")

    # Budget so right panel is visible when present.
    full_budget = int(max_chars * 0.40)
    side_budget = int(max_chars * 0.28)
    right_budget = max_chars - (len("[Page 1]\n[Region: full]\n") + full_budget + len("\n[Region: left]\n") + side_budget + len("\n[Region: right]\n"))
    right_budget = max(80, min(int(max_chars * 0.32), right_budget))

    parts = ["[Page 1]"]
    if full:
        parts.extend(["[Region: full]", full[:full_budget].rstrip()])
    if left:
        parts.extend(["[Region: left]", left[:side_budget].rstrip()])
    if right:
        parts.extend(["[Region: right]", right[:right_budget].rstrip()])

    composed = "\n".join(p for p in parts if p).strip()
    if len(composed) > max_chars:
        composed = composed[:max_chars].rstrip()
    return composed


class TextExtractor:
    """Service for extracting text from various document formats"""
    
    def __init__(self, ocr_service=None):
        """
        Initialize text extractor.
        
        Args:
            ocr_service: Optional OCRService instance for scanned document processing
        """
        self.supported_extensions = set(BASE_SUPPORTED_EXTENSIONS)
        if ocr_service and ocr_service.is_available():
            self.supported_extensions.update(IMAGE_EXTENSIONS_OCR)
        
        self.ocr_service = ocr_service
        
        # Excel processing limits (configurable for performance)
        self.max_sheets_per_file = 20  # Limit sheets to prevent memory issues
        self.max_rows_per_sheet = 10000  # Limit rows per sheet
        # PowerPoint processing limits
        self.max_slides_per_presentation = 200  # Limit slides to prevent memory issues
    
    async def extract_text_async(self, file_path: str) -> str:
        """Extract text asynchronously to avoid blocking"""
        def sync_extract():
            return self._extract_text_sync(file_path)
        
        # Run in thread pool for CPU-intensive operations
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, sync_extract)
    
    def _extract_text_sync(self, file_path: str) -> str:
        """Synchronous text extraction with better error handling"""
        path_obj = Path(file_path)
        ext = path_obj.suffix.lower()
        
        try:
            if ext == ".pdf":
                return self._extract_pdf(file_path)
            elif ext == ".docx":
                return self._extract_docx(file_path)
            elif ext == ".txt":
                return self._extract_txt(file_path)
            elif ext in {".xlsx", ".xls"}:
                return self._extract_excel(file_path)
            elif ext == ".pptx":
                return self._extract_pptx(file_path)
            elif ext in {".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp"}:
                return self._extract_image(file_path)
            else:
                raise ValueError(f"Unsupported file type: {ext}")
        except Exception as e:
            logger.error(f"Text extraction failed for {file_path}: {e}")
            raise
    
    def _extract_pdf(self, file_path: str) -> str:
        """
        Extract text from PDF.
        If PDF is scanned (no extractable text), uses OCR if available.
        """
        try:
            import fitz
            doc = fitz.open(file_path)
            text_parts = []
            for page_idx, page in enumerate(doc, start=1):
                # Industry-standard layout reconstruction:
                # - Use "dict" (lines/spans with bounding boxes), not "blocks".
                #   Blocks are often merged across columns/boxes, which is exactly
                #   what caused DELIVER TO + REFERENCES to flatten together.
                page_rect = page.rect
                page_w = float(page_rect.width) if page_rect else 0.0
                page_h = float(page_rect.height) if page_rect else 0.0
                x_mid = page_w / 2.0 if page_w else 0.0

                extracted = page.get_text("dict") or {}
                blocks = extracted.get("blocks") or []

                full_lines = []
                left_lines = []
                right_lines = []

                def _is_full_width(x0f: float, x1f: float) -> bool:
                    if not page_w:
                        return False
                    width = max(0.0, x1f - x0f)
                    return bool(
                        width >= 0.55 * page_w
                        or (x0f <= 0.15 * page_w and x1f >= 0.85 * page_w)
                    )

                for b in blocks:
                    if not isinstance(b, dict):
                        continue
                    lines = b.get("lines") or []
                    for ln in lines:
                        if not isinstance(ln, dict):
                            continue
                        bbox = ln.get("bbox") or None
                        spans = ln.get("spans") or []
                        if not bbox or len(bbox) != 4 or not spans:
                            continue
                        x0, y0, x1, y1 = bbox
                        try:
                            x0f, y0f, x1f, y1f = float(x0), float(y0), float(x1), float(y1)
                        except (TypeError, ValueError):
                            continue

                        # Reconstruct line text by x-order spans.
                        span_texts = []
                        span_items = []
                        for sp in spans:
                            if not isinstance(sp, dict):
                                continue
                            t = sp.get("text") or ""
                            sb = sp.get("bbox") or None
                            if not t or not isinstance(t, str):
                                continue
                            if sb and len(sb) == 4:
                                try:
                                    sx0 = float(sb[0])
                                except (TypeError, ValueError):
                                    sx0 = x0f
                            else:
                                sx0 = x0f
                            span_items.append((sx0, t))
                        if not span_items:
                            continue
                        span_items.sort(key=lambda t: t[0])
                        line_text = "".join(t[1] for t in span_items).strip()
                        if not line_text:
                            continue

                        entry = (y0f, x0f, line_text)
                        if _is_full_width(x0f, x1f):
                            full_lines.append(entry)
                            continue
                        cx = (x0f + x1f) / 2.0
                        if x_mid and cx < x_mid:
                            left_lines.append(entry)
                        else:
                            right_lines.append(entry)

                # If dict extraction produced nothing (rare), fall back to plain text.
                if not (full_lines or left_lines or right_lines):
                    t = page.get_text() or ""
                    if t.strip():
                        text_parts.append(f"[Page {page_idx}]\n{t.strip()}")
                    continue

                full_lines.sort(key=lambda t: (t[0], t[1]))
                left_lines.sort(key=lambda t: (t[0], t[1]))
                right_lines.sort(key=lambda t: (t[0], t[1]))

                # Keep region labels, but preserve vertical reading intent:
                # - "full" headers at the top should appear first
                # - then left/right panels
                # - then remaining full-width body text (e.g. long cargo description)
                full_top = []
                full_rest = []
                # Heuristic: treat "top of page" full-width lines as headers/titles.
                # This is more stable than comparing against panel y because PDF
                # extraction sometimes reports body blocks with unexpectedly small y.
                if not page_h:
                    full_top = full_lines
                else:
                    top_cut = 0.32 * page_h
                    for ln in full_lines:
                        (full_top if ln[0] <= top_cut else full_rest).append(ln)

                page_lines = [f"[Page {page_idx}]"]
                if full_top:
                    page_lines.append("[Region: full]")
                    page_lines.append("\n".join(t[2] for t in full_top))
                if left_lines:
                    page_lines.append("[Region: left]")
                    page_lines.append("\n".join(t[2] for t in left_lines))
                if right_lines:
                    page_lines.append("[Region: right]")
                    page_lines.append("\n".join(t[2] for t in right_lines))
                if full_rest:
                    page_lines.append("[Region: full]")
                    page_lines.append("\n".join(t[2] for t in full_rest))
                page_text = "\n".join(page_lines).strip()
                if page_text:
                    text_parts.append(page_text)
            doc.close()
            text = "\n\n".join(text_parts)
            
            # Check if PDF has extractable text
            if text.strip():
                # PDF has text, return it
                return text
            else:
                # PDF appears to be scanned (no text), try OCR
                if self.ocr_service and self.ocr_service.is_available():
                    logger.info(f"PDF {file_path} has no extractable text, attempting OCR...")
                    try:
                        # Use sync method directly (OCR service has sync methods for internal use)
                        ocr_text = self.ocr_service._extract_text_from_scanned_pdf_sync(file_path)
                        return ocr_text if ocr_text else ""
                    except Exception as ocr_error:
                        logger.warning(f"OCR failed for PDF {file_path}: {ocr_error}")
                        return ""
                else:
                    logger.warning(f"PDF {file_path} has no extractable text and OCR is not available")
                    return ""
                    
        except Exception as e:
            logger.error(f"PDF extraction failed for {file_path}: {e}")
            raise
    
    def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX"""
        try:
            import docx
            doc = docx.Document(file_path)
            parts = [para.text for para in doc.paragraphs]
            for table in doc.tables:
                for row in table.rows:
                    parts.append("\t".join(cell.text for cell in row.cells))
            return "\n".join(parts)
        except Exception as e:
            logger.error(f"DOCX extraction failed for {file_path}: {e}")
            raise
    
    def _extract_pptx(self, file_path: str) -> str:
        """
        Extract text from PowerPoint files (.pptx).
        
        Handles:
        - Slide titles and content
        - Text boxes and paragraphs
        - Tables
        - Multiple slides with slide numbers
        
        Args:
            file_path: Path to PowerPoint file
            
        Returns:
            Structured text representation of PowerPoint content
        """
        try:
            from pptx import Presentation
            
            # Open presentation with error handling
            try:
                prs = Presentation(file_path)
            except Exception as e:
                logger.error(f"Failed to open PowerPoint file {file_path}: {e}")
                raise
            
            text_parts = []
            
            # Safely get slides count
            try:
                slides_list = list(prs.slides) if hasattr(prs, 'slides') else []
                slide_count = min(len(slides_list), self.max_slides_per_presentation)
                
                logger.info(f"PowerPoint file {file_path}: Found {len(slides_list)} slides, processing {slide_count}")
                
                if len(slides_list) > self.max_slides_per_presentation:
                    logger.warning(
                        f"PowerPoint file {file_path} has {len(slides_list)} slides, "
                        f"processing only first {self.max_slides_per_presentation}"
                    )
            except Exception as e:
                logger.error(f"Failed to access slides in {file_path}: {e}")
                return ""
            
            for slide_idx, slide in enumerate(slides_list[:slide_count], start=1):
                try:
                    slide_text_parts = [f"Slide {slide_idx}:"]
                    
                    # Extract text from all shapes
                    # Use minimal approach to avoid 'rId' errors with problematic PowerPoint files
                    try:
                        # Try to get shapes - this might fail with rId error
                        try:
                            shapes_iter = iter(slide.shapes) if hasattr(slide, 'shapes') else iter([])
                            shapes = list(shapes_iter)
                        except AttributeError as e:
                            if "'list' object has no attribute 'rId'" in str(e) or "rId" in str(e):
                                logger.warning(f"PowerPoint file has rId error on slide {slide_idx}, skipping shape iteration")
                                shapes = []
                            else:
                                raise
                        
                        for shape in shapes:
                            try:
                                # Try text_frame first (most reliable)
                                text_extracted = False
                                if hasattr(shape, "has_text_frame"):
                                    try:
                                        if shape.has_text_frame:
                                            text_frame = shape.text_frame
                                            if text_frame:
                                                # Extract text from paragraphs and runs (more comprehensive)
                                                para_texts = []
                                                try:
                                                    # Method 1: Extract from paragraphs (standard approach)
                                                    for paragraph in text_frame.paragraphs:
                                                        para_text = paragraph.text
                                                        if para_text and para_text.strip():
                                                            para_texts.append(para_text.strip())
                                                    
                                                    # Method 2: Also extract from runs (handles formatted text better)
                                                    # This is important for Canva files which may have complex formatting
                                                    if not para_texts:
                                                        for paragraph in text_frame.paragraphs:
                                                            run_texts = []
                                                            for run in paragraph.runs:
                                                                if run.text and run.text.strip():
                                                                    run_texts.append(run.text.strip())
                                                            if run_texts:
                                                                para_texts.append(" ".join(run_texts))
                                                except Exception as para_e:
                                                    # If paragraph iteration fails, try direct text_frame.text
                                                    try:
                                                        direct_text = text_frame.text
                                                        if direct_text and direct_text.strip():
                                                            para_texts.append(direct_text.strip())
                                                    except Exception:
                                                        logger.debug(f"Error extracting from text_frame on slide {slide_idx}: {para_e}")
                                                
                                                if para_texts:
                                                    slide_text_parts.append("\n".join(para_texts))
                                                    text_extracted = True
                                    except (AttributeError, TypeError) as e:
                                        if "rId" in str(e):
                                            logger.debug(f"Skipping shape with rId error on slide {slide_idx}")
                                            continue
                                        logger.debug(f"Error reading text_frame: {e}")
                                
                                # Fallback: try direct text access (less reliable but sometimes works)
                                if not text_extracted and hasattr(shape, "text"):
                                    try:
                                        text_content = shape.text
                                        if text_content and text_content.strip():
                                            slide_text_parts.append(text_content.strip())
                                            text_extracted = True
                                    except (AttributeError, TypeError) as e:
                                        if "rId" in str(e):
                                            logger.debug(f"Skipping shape with rId error on slide {slide_idx}")
                                            continue
                                        logger.debug(f"Error reading text: {e}")
                                
                                # Handle grouped shapes (Canva often uses grouped shapes)
                                if not text_extracted and hasattr(shape, "shapes"):
                                    try:
                                        # This is a group shape - recursively extract from nested shapes
                                        group_texts = []
                                        for sub_shape in shape.shapes:
                                            if hasattr(sub_shape, "has_text_frame") and sub_shape.has_text_frame:
                                                try:
                                                    sub_text = sub_shape.text_frame.text
                                                    if sub_text and sub_text.strip():
                                                        group_texts.append(sub_text.strip())
                                                except Exception:
                                                    pass
                                        if group_texts:
                                            slide_text_parts.append("\n".join(group_texts))
                                            text_extracted = True
                                    except Exception as e:
                                        logger.debug(f"Error processing group shape on slide {slide_idx}: {e}")
                                
                                # Extract text from tables (if no rId error occurred)
                                if hasattr(shape, "has_table"):
                                    try:
                                        if shape.has_table and hasattr(shape, "table"):
                                            table = shape.table
                                            table_text = self._extract_table_from_slide(table)
                                            if table_text:
                                                slide_text_parts.append(f"Table:\n{table_text}")
                                    except (AttributeError, TypeError) as e:
                                        if "rId" in str(e):
                                            logger.debug(f"Skipping table with rId error on slide {slide_idx}")
                                            continue
                                        logger.debug(f"Error extracting table: {e}")
                            except (AttributeError, TypeError) as e:
                                if "rId" in str(e):
                                    logger.debug(f"Skipping shape with rId error on slide {slide_idx}")
                                    continue
                                logger.debug(f"Error processing shape on slide {slide_idx}: {e}")
                                continue
                    except (AttributeError, TypeError) as e:
                        if "rId" in str(e):
                            logger.warning(f"PowerPoint file has rId error accessing shapes on slide {slide_idx}")
                        else:
                            logger.warning(f"Error processing shapes on slide {slide_idx}: {e}")
                    
                    # Only add slides with content
                    if len(slide_text_parts) > 1:  # More than just "Slide X:"
                        slide_content = "\n".join(slide_text_parts)
                        text_parts.append(slide_content)
                        logger.debug(f"Slide {slide_idx}: Extracted {len(slide_content)} characters")
                    else:
                        logger.debug(f"Slide {slide_idx}: No extractable content")
                except Exception as e:
                    # If entire slide processing fails, log and continue with next slide
                    logger.warning(f"Failed to process slide {slide_idx} in {file_path}: {e}")
                    continue
            
            if not text_parts:
                logger.warning(f"No extractable content found in {file_path}")
                return ""
            
            extracted_text = "\n\n".join(text_parts)
            logger.info(f"PowerPoint extraction complete for {file_path}: {len(slides_list)} slides processed, {len(text_parts)} slides with content, {len(extracted_text)} total characters")
            return extracted_text
            
        except Exception as e:
            error_msg = str(e).lower()
            if "password" in error_msg or "encrypted" in error_msg:
                raise ValueError(f"PowerPoint file is password-protected: {file_path}")
            elif "corrupt" in error_msg or "invalid" in error_msg:
                raise ValueError(f"PowerPoint file appears to be corrupted: {file_path}")
            else:
                logger.error(f"PPTX extraction failed for {file_path}: {e}")
                raise
    
    def _extract_table_from_slide(self, table) -> str:
        """
        Extract text from a PowerPoint table.
        
        Args:
            table: PowerPoint table object
            
        Returns:
            Formatted table text with headers and rows
        """
        try:
            lines = []
            
            # Safely access table rows
            if not hasattr(table, 'rows'):
                return ""
            
            # Convert rows to list to avoid slice/int comparison errors
            try:
                rows_list = list(table.rows) if table.rows else []
            except (TypeError, AttributeError) as e:
                logger.debug(f"Error accessing table rows: {e}")
                return ""
            
            if not rows_list:
                return ""
            
            # Extract headers (first row)
            try:
                header_row = rows_list[0]
                headers = []
                if hasattr(header_row, 'cells'):
                    for cell in header_row.cells:
                        try:
                            cell_text = cell.text.strip() if cell.text else ""
                            headers.append(cell_text)
                        except Exception as e:
                            logger.debug(f"Error reading header cell: {e}")
                            headers.append("")
                
                if headers:
                    lines.append(f"Headers: {' | '.join(headers)}")
            except (IndexError, AttributeError) as e:
                logger.debug(f"Error extracting table headers: {e}")
            
            # Extract data rows (skip first row as it's the header)
            try:
                for row_idx, row in enumerate(rows_list[1:], start=1):
                    try:
                        row_values = []
                        has_content = False
                        
                        if hasattr(row, 'cells'):
                            for cell in row.cells:
                                try:
                                    cell_text = cell.text.strip() if cell.text else ""
                                    row_values.append(cell_text)
                                    if cell_text:
                                        has_content = True
                                except Exception as e:
                                    logger.debug(f"Error reading cell: {e}")
                                    row_values.append("")
                        
                        if has_content:
                            lines.append(f"Row {row_idx}: {' | '.join(row_values)}")
                    except Exception as e:
                        logger.debug(f"Error processing table row {row_idx}: {e}")
                        continue
            except Exception as e:
                logger.debug(f"Error iterating table rows: {e}")
            
            return "\n".join(lines) if lines else ""
            
        except Exception as e:
            logger.warning(f"Failed to extract table from slide: {e}")
            return ""
    
    def _extract_txt(self, file_path: str) -> str:
        """Extract text from TXT with encoding detection"""
        try:
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            for encoding in encodings:
                try:
                    with open(file_path, "r", encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            
            # Fallback with error handling
            with open(file_path, "r", encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            logger.error(f"TXT extraction failed for {file_path}: {e}")
            raise
    
    def _extract_excel(self, file_path: str) -> str:
        """
        Extract text from Excel files (.xlsx and .xls).
        
        Handles:
        - Multiple sheets (with sheet names as context)
        - Table structure preservation
        - Headers detection
        - Empty cells gracefully
        - Large files with configurable limits
        
        Args:
            file_path: Path to Excel file
            
        Returns:
            Structured text representation of Excel content
        """
        path_obj = Path(file_path)
        ext = path_obj.suffix.lower()
        
        if ext == ".xlsx":
            return self._extract_xlsx(file_path)
        elif ext == ".xls":
            return self._extract_xls(file_path)
        else:
            raise ValueError(f"Unsupported Excel format: {ext}")
    
    def _extract_xlsx(self, file_path: str) -> str:
        """Extract text from .xlsx files using openpyxl"""
        try:
            from openpyxl import load_workbook
            
            workbook = load_workbook(file_path, read_only=True, data_only=True)
            text_parts = []
            
            # Process each sheet (with limit to prevent memory issues)
            sheet_names = workbook.sheetnames[:self.max_sheets_per_file]
            
            if len(workbook.sheetnames) > self.max_sheets_per_file:
                logger.warning(
                    f"Excel file {file_path} has {len(workbook.sheetnames)} sheets, "
                    f"processing only first {self.max_sheets_per_file}"
                )
            
            for sheet_name in sheet_names:
                sheet = workbook[sheet_name]
                sheet_text = self._extract_sheet_content(sheet, sheet_name)
                if sheet_text.strip():  # Only add non-empty sheets
                    text_parts.append(sheet_text)
            
            workbook.close()
            
            if not text_parts:
                logger.warning(f"No extractable content found in {file_path}")
                return ""
            
            return "\n\n".join(text_parts)
            
        except Exception as e:
            error_msg = str(e).lower()
            if "password" in error_msg or "encrypted" in error_msg:
                raise ValueError(f"Excel file is password-protected: {file_path}")
            elif "corrupt" in error_msg or "invalid" in error_msg:
                raise ValueError(f"Excel file appears to be corrupted: {file_path}")
            else:
                logger.error(f"XLSX extraction failed for {file_path}: {e}")
                raise
    
    def _extract_xls(self, file_path: str) -> str:
        """Extract text from legacy .xls files using xlrd"""
        try:
            import xlrd
            
            workbook = xlrd.open_workbook(file_path, on_demand=True)
            text_parts = []
            
            # Process each sheet (with limit)
            sheet_count = min(len(workbook.sheet_names()), self.max_sheets_per_file)
            
            if len(workbook.sheet_names()) > self.max_sheets_per_file:
                logger.warning(
                    f"Excel file {file_path} has {len(workbook.sheet_names())} sheets, "
                    f"processing only first {self.max_sheets_per_file}"
                )
            
            for sheet_idx in range(sheet_count):
                sheet = workbook.sheet_by_index(sheet_idx)
                sheet_name = sheet.name
                sheet_text = self._extract_xls_sheet_content(sheet, sheet_name)
                if sheet_text.strip():  # Only add non-empty sheets
                    text_parts.append(sheet_text)
            
            workbook.release_resources()
            
            if not text_parts:
                logger.warning(f"No extractable content found in {file_path}")
                return ""
            
            return "\n\n".join(text_parts)
            
        except xlrd.biffh.XLRDError as e:
            error_msg = str(e).lower()
            if "password" in error_msg or "encrypted" in error_msg:
                raise ValueError(f"Excel file is password-protected: {file_path}")
            elif "corrupt" in error_msg:
                raise ValueError(f"Excel file appears to be corrupted: {file_path}")
            else:
                logger.error(f"XLS extraction failed for {file_path}: {e}")
                raise
        except Exception as e:
            logger.error(f"XLS extraction failed for {file_path}: {e}")
            raise
    
    def _extract_sheet_content(self, sheet, sheet_name: str) -> str:
        """
        Extract content from an openpyxl sheet.
        
        Formats output as structured text with headers and rows.
        """
        lines = [f"Sheet: {sheet_name}"]
        
        # Get used range
        if sheet.max_row == 0 or sheet.max_column == 0:
            return ""  # Empty sheet
        
        # Limit rows to prevent memory issues
        max_row = min(sheet.max_row, self.max_rows_per_sheet)
        
        if sheet.max_row > self.max_rows_per_sheet:
            logger.debug(
                f"Sheet '{sheet_name}' has {sheet.max_row} rows, "
                f"processing only first {self.max_rows_per_sheet}"
            )
        
        # Extract headers (first row)
        headers = []
        for col_idx in range(1, min(sheet.max_column + 1, 100)):  # Limit columns too
            cell_value = sheet.cell(row=1, column=col_idx).value
            if cell_value is not None:
                headers.append(str(cell_value).strip())
            else:
                headers.append("")
        
        # Remove trailing empty headers
        while headers and not headers[-1]:
            headers.pop()
        
        if headers:
            lines.append(f"Headers: {' | '.join(headers)}")
        
        # Extract data rows
        row_count = 0
        for row_idx in range(2, max_row + 1):  # Start from row 2 (skip header)
            row_values = []
            has_content = False
            
            for col_idx in range(1, len(headers) + 1):
                cell_value = sheet.cell(row=row_idx, column=col_idx).value
                if cell_value is not None:
                    # Convert to string, handling different types
                    if isinstance(cell_value, (int, float)):
                        cell_str = str(cell_value)
                    elif isinstance(cell_value, str):
                        cell_str = cell_value.strip()
                    else:
                        cell_str = str(cell_value)
                    row_values.append(cell_str)
                    if cell_str:
                        has_content = True
                else:
                    row_values.append("")
            
            # Only add rows with content
            if has_content:
                lines.append(f"Row {row_idx - 1}: {' | '.join(row_values)}")
                row_count += 1
        
        if row_count == 0:
            return ""  # Sheet has no data rows
        
        return "\n".join(lines)
    
    def _extract_xls_sheet_content(self, sheet, sheet_name: str) -> str:
        """
        Extract content from an xlrd sheet (legacy .xls format).
        
        Formats output as structured text with headers and rows.
        """
        lines = [f"Sheet: {sheet_name}"]
        
        if sheet.nrows == 0 or sheet.ncols == 0:
            return ""  # Empty sheet
        
        # Limit rows
        max_row = min(sheet.nrows, self.max_rows_per_sheet)
        
        if sheet.nrows > self.max_rows_per_sheet:
            logger.debug(
                f"Sheet '{sheet_name}' has {sheet.nrows} rows, "
                f"processing only first {self.max_rows_per_sheet}"
            )
        
        # Extract headers (first row)
        headers = []
        for col_idx in range(min(sheet.ncols, 100)):  # Limit columns
            cell_value = sheet.cell_value(0, col_idx)
            if cell_value:
                headers.append(str(cell_value).strip())
            else:
                headers.append("")
        
        # Remove trailing empty headers
        while headers and not headers[-1]:
            headers.pop()
        
        if headers:
            lines.append(f"Headers: {' | '.join(headers)}")
        
        # Extract data rows
        row_count = 0
        for row_idx in range(1, max_row):  # Start from row 1 (skip header)
            row_values = []
            has_content = False
            
            for col_idx in range(len(headers)):
                if col_idx < sheet.ncols:
                    cell_value = sheet.cell_value(row_idx, col_idx)
                    if cell_value:
                        cell_str = str(cell_value).strip()
                        row_values.append(cell_str)
                        if cell_str:
                            has_content = True
                    else:
                        row_values.append("")
                else:
                    row_values.append("")
            
            # Only add rows with content
            if has_content:
                lines.append(f"Row {row_idx}: {' | '.join(row_values)}")
                row_count += 1
        
        if row_count == 0:
            return ""  # Sheet has no data rows
        
        return "\n".join(lines)
    
    def _extract_image(self, file_path: str) -> str:
        """
        Extract text from image file using OCR.
        
        Args:
            file_path: Path to image file
            
        Returns:
            Extracted text from image
        """
        if not self.ocr_service or not self.ocr_service.is_available():
            raise RuntimeError(
                f"OCR not available. Cannot extract text from image file: {file_path}. "
                "Please install Tesseract OCR to enable image processing."
            )
        
        try:
            # Use sync method directly (OCR service has sync methods for internal use)
            text = self.ocr_service._extract_text_from_image_sync(file_path)
            
            if not text:
                logger.warning(f"No text extracted from image {file_path}")
            
            return text
            
        except Exception as e:
            logger.error(f"Image OCR extraction failed for {file_path}: {e}")
            raise
    
    def is_supported_file(self, file_path: str) -> bool:
        """Check if file type is supported"""
        return Path(file_path).suffix.lower() in self.supported_extensions
